"""Ingest the D:\\University coursework tree into the vault as linked notes.

Pull-only today (filesearch); this module turns the tree into markdown under
<vault>/coursework/ so it joins the semantic index and the Obsidian graph.
The walk is strictly read-only. Called by voice/heartbeat.py's
_check_university_intake task.
"""
from __future__ import annotations

import hashlib
import re
import sys
import unicodedata
from dataclasses import dataclass, field
from datetime import date, timedelta
from pathlib import Path

_KINDS = {"uni assignments": "assignment", "uni tutorial": "tutorial"}
_PROGRAMS = {"degree", "diploma"}


@dataclass(frozen=True)
class ParsedPath:
    category: str      # "assignment" | "tutorial"
    program: str       # "Degree" | "Diploma"
    semester: int
    course: str        # subject code, e.g. "BIT216"
    subpath: str       # posix relpath below the course folder
    source: Path       # absolute original path


def parse_path(p: Path, root: Path) -> ParsedPath | None:
    """Map an absolute file path to a ParsedPath, or None if it does not match
    `<Uni Assignments|Uni Tutorial>/<program>/Sem <N>/<course>/<...file>`
    anywhere below `root`."""
    try:
        rel = p.relative_to(root)
    except ValueError:
        return None
    parts = list(rel.parts)
    kind_ix = next(
        (i for i, seg in enumerate(parts) if seg.lower() in _KINDS), None
    )
    if kind_ix is None or len(parts) < kind_ix + 5:
        return None
    program = parts[kind_ix + 1]
    sem_seg = parts[kind_ix + 2]
    course = parts[kind_ix + 3]
    if program.lower() not in _PROGRAMS:
        return None
    if not sem_seg.lower().startswith("sem "):
        return None
    try:
        semester = int(sem_seg.split(None, 1)[1])
    except (IndexError, ValueError):
        return None
    subpath = "/".join(parts[kind_ix + 4:])
    if not subpath:
        return None
    return ParsedPath(
        category=_KINDS[parts[kind_ix].lower()],
        program=program,
        semester=semester,
        course=course,
        subpath=subpath,
        source=p,
    )


def is_denied(p: Path, denylist: list[str], allow_ext: list[str],
              *, _denied: set[str] | None = None,
              _allow: set[str] | None = None) -> bool:
    """True if any path segment is in `denylist` (case-insensitive) or the
    file extension is not in `allow_ext`. `_denied`/`_allow` let callers pass
    pre-lowercased sets so they are not rebuilt per file."""
    denied = _denied if _denied is not None else {d.lower() for d in denylist}
    allow = _allow if _allow is not None else {e.lower() for e in allow_ext}
    if any(seg.lower() in denied for seg in p.parts):
        return True
    return p.suffix.lower() not in allow


def _extract_txt(p: Path) -> str:
    return p.read_text(encoding="utf-8", errors="replace")


def _extract_pdf(p: Path) -> str:
    from pypdf import PdfReader  # lazy: missing lib -> ImportError caught by extract_text
    reader = PdfReader(str(p))
    return "\n".join((page.extract_text() or "") for page in reader.pages)


def _extract_docx(p: Path) -> str:
    import docx  # python-docx
    d = docx.Document(str(p))
    return "\n".join(par.text for par in d.paragraphs)


def _extract_pptx(p: Path) -> str:
    import pptx  # python-pptx
    prs = pptx.Presentation(str(p))
    lines: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                lines.append(shape.text_frame.text)
    return "\n".join(lines)


_EXTRACTOR_NAMES = {
    ".pdf": "_extract_pdf", ".docx": "_extract_docx", ".pptx": "_extract_pptx",
    ".txt": "_extract_txt", ".md": "_extract_txt",
}


def extract_text(p: Path) -> tuple[str, bool, str | None]:
    """(text, ok, error). ok is True only when an extractor ran and produced
    non-empty text. error is the exception repr when an extractor ran and
    raised (corrupt file), else None -- this lets the caller tell a corrupt
    file apart from a genuinely empty / image-only one. Missing lib / parse
    error / unsupported extension -> ("", False, ...) so the caller can still
    write a stub note. Resolves the extractor late so tests can monkeypatch
    _extract_*."""
    name = _EXTRACTOR_NAMES.get(p.suffix.lower())
    if name is None:
        return "", False, None
    fn = getattr(sys.modules[__name__], name)
    try:
        text = (fn(p) or "").strip()
    except Exception as e:
        return "", False, repr(e)
    return (text, True, None) if text else ("", False, None)


def slugify(subpath: str) -> str:
    stem = subpath.rsplit(".", 1)[0]
    norm = unicodedata.normalize("NFKD", stem).encode("ascii", "ignore").decode()
    return re.sub(r"-+", "-", re.sub(r"[^a-z0-9]+", "-", norm.lower())).strip("-")


def summarise(text: str, parsed: ParsedPath) -> str:
    words = text.split()
    if words:
        return " ".join(words[:60])
    return (
        f"{parsed.category.capitalize()} for {parsed.course} "
        f"({parsed.program} Sem {parsed.semester}). "
        f"Source file: {parsed.subpath.rsplit('/', 1)[-1]}. No extractable text."
    )


def _frontmatter(parsed: ParsedPath, truncated: bool) -> str:
    return (
        "---\n"
        "type: coursework\n"
        f"course: {parsed.course}\n"
        f"program: {parsed.program}\n"
        f"semester: {parsed.semester}\n"
        f"category: {parsed.category}\n"
        f"source: {parsed.source.as_posix()}\n"
        f"added: {date.today().isoformat()}\n"
        f"truncated: {'true' if truncated else 'false'}\n"
        "---\n"
    )


def note_for(parsed: ParsedPath, text: str, truncated: bool, mode: str) -> tuple[str, str]:
    course_l = parsed.course.lower()
    stem = parsed.subpath.rsplit("/", 1)[-1].rsplit(".", 1)[0]
    if mode == "rollup":
        relpath = f"coursework/{course_l}.md"
        md = (
            f"## {date.today().isoformat()} — {parsed.subpath}\n\n"
            f"{summarise(text, parsed)}\n\n"
            f"{text}\n"
        )
        return relpath, md
    relpath = f"coursework/{course_l}/{slugify(parsed.subpath)}.md"
    md = (
        f"{_frontmatter(parsed, truncated)}\n"
        f"# {parsed.course} — {parsed.category.capitalize()}: {stem}\n\n"
        f"{summarise(text, parsed)}\n\n"
        f"[[{course_l}]]\n\n"
        f"---\n\n"
        f"{text}"
    )
    return relpath, md


_MONTHS = {m.lower(): i for i, m in enumerate(
    ["January", "February", "March", "April", "May", "June", "July",
     "August", "September", "October", "November", "December"], start=1)}

_DATE_RES = [
    ("iso", re.compile(r"\b(\d{4})-(\d{2})-(\d{2})\b")),                # 2026-10-15
    ("dmy", re.compile(r"\b(\d{1,2})[/-](\d{1,2})[/-](\d{4})\b")),      # 15/10/2026
    ("dMy", re.compile(r"\b(\d{1,2})\s+([A-Za-z]+)\s+(\d{4})\b")),      # 5 December 2026
    ("Mdy", re.compile(r"\b([A-Za-z]+)\s+(\d{1,2}),?\s+(\d{4})\b")),    # December 5, 2026
]
_KW_RE = re.compile(r"\b(due|deadline|submission|submit|by)\b", re.IGNORECASE)


def _iso_for(kind: str, g: tuple[str, ...]) -> str | None:
    try:
        if kind == "iso":
            y, mo, d = int(g[0]), int(g[1]), int(g[2])
        elif kind == "dmy":
            d, mo, y = int(g[0]), int(g[1]), int(g[2])
        elif kind == "dMy":
            d, mo, y = int(g[0]), _MONTHS.get(g[1].lower(), 0), int(g[2])
        else:  # "Mdy"
            mo, d, y = _MONTHS.get(g[0].lower(), 0), int(g[1]), int(g[2])
        if not mo:
            return None
        return date(y, mo, d).isoformat()
    except ValueError:
        return None


def detect_deadlines(text: str, parsed: ParsedPath,
                     today: date | None = None) -> list[tuple[str, str]]:
    if parsed.category != "assignment":
        return []
    today = today or date.today()
    horizon = today + timedelta(days=180)
    window = text[:4000]
    stem = parsed.subpath.rsplit("/", 1)[-1].rsplit(".", 1)[0]
    title = f"{parsed.course} — {stem}"
    cands: list[tuple[int, str]] = []  # (document position, iso)
    for kind, rex in _DATE_RES:
        for m in rex.finditer(window):
            ctx = window[max(0, m.start() - 40): m.end() + 40]
            if not _KW_RE.search(ctx):
                continue
            iso = _iso_for(kind, m.groups())
            if not iso:
                continue
            d = date.fromisoformat(iso)
            if d < today or d > horizon:
                continue
            cands.append((m.start(), iso))

    seen: set[str] = set()
    out: list[tuple[str, str]] = []
    for _pos, iso in sorted(cands):
        if iso in seen:
            continue
        seen.add(iso)
        out.append((iso, title))
        if len(out) == 2:
            break
    return out


@dataclass
class IntakeResult:
    added: list[str] = field(default_factory=list)
    updated: list[str] = field(default_factory=list)
    skipped: int = 0
    unparsed: int = 0
    errors: list[str] = field(default_factory=list)
    deadlines: list[tuple[str, str]] = field(default_factory=list)
    partial: bool = False   # per-tick write cap hit; remainder resumes next tick


def file_hash(p: Path) -> str:
    return hashlib.sha256(p.read_bytes()).hexdigest()[:16]


def run_intake(root: Path, vault: Path, *, manifest: dict, config: dict) -> IntakeResult:
    res = IntakeResult()
    if not root.is_dir():
        return res
    manifest.setdefault("version", 1)
    files: dict = manifest.setdefault("files", {})

    exts = config["university_intake_extensions"]
    denylist = config["university_intake_denylist"]
    mode = config.get("university_intake_mode", "per-file")
    cap = int(config.get("university_intake_max_chars", 20000))
    want_deadlines = bool(config.get("university_intake_deadline_detection", True))
    max_per_tick = int(config.get("university_intake_max_files_per_tick", 200))

    denied_set = {d.lower() for d in denylist}
    allow_set = {e.lower() for e in exts}

    # note-relpath -> source relposix that owns it, so a slug collision inside
    # one course folder (brief.docx + brief.pdf) is deduped deterministically
    # and stably across runs (spec §6). Pre-seeded from the manifest so every
    # already-owned note path is reserved before the walk starts -- a colliding
    # sibling that sorts *earlier* than its owner must not clobber it.
    claimed: dict[str, str] = {}
    if mode != "rollup":
        for _r, _m in files.items():
            if _m.get("note"):
                claimed.setdefault(_m["note"], _r)

    def _claim(relnote: str, src_rel: str) -> str:
        owner = claimed.get(relnote)
        if owner is None or owner == src_rel:
            claimed[relnote] = src_rel
            return relnote
        stem = relnote[:-3] if relnote.endswith(".md") else relnote
        n = 2
        while True:
            cand = f"{stem}-{n}.md"
            owner = claimed.get(cand)
            if owner is None or owner == src_rel:
                claimed[cand] = src_rel
                return cand
            n += 1

    written = 0
    seen_rel: set[str] = set()
    for p in sorted(root.rglob("*")):
        if not p.is_file():
            continue
        if is_denied(p, denylist, exts, _denied=denied_set, _allow=allow_set):
            res.skipped += 1
            continue
        parsed = parse_path(p, root)
        if parsed is None:
            res.unparsed += 1
            continue
        rel = p.relative_to(root).as_posix()
        seen_rel.add(rel)
        prior = files.get(rel)
        if mode != "rollup" and prior and prior.get("note"):
            claimed.setdefault(prior["note"], rel)
        try:
            st = p.stat()
            # Stat short-circuit: same size + mtime as the manifest recorded
            # => unchanged, skip without reading the file (spec §4.3).
            if (prior and prior.get("mtime") == st.st_mtime
                    and prior.get("size") == st.st_size):
                res.skipped += 1
                continue

            h = file_hash(p)
            if prior and prior.get("hash") == h:
                # Content unchanged; refresh the stat cache so the next tick
                # short-circuits before hashing.
                prior["mtime"] = st.st_mtime
                prior["size"] = st.st_size
                res.skipped += 1
                continue

            if written >= max_per_tick:
                # Leave in seen_rel (not pruned) with no manifest write so the
                # next tick resumes here (spec §11 backfill throttle).
                res.partial = True
                continue

            text, _ok, err = extract_text(p)
            if err is not None and len(res.errors) < 20:
                res.errors.append(f"{rel}: extract failed: {err}")
            truncated = len(text) > cap
            body = text[:cap] if truncated else text
            relnote, md = note_for(parsed, body, truncated, mode)
            if mode != "rollup":
                relnote = (
                    prior["note"] if prior and prior.get("note")
                    else _claim(relnote, rel)
                )
            dest = vault / relnote
            dest.parent.mkdir(parents=True, exist_ok=True)
            if mode == "rollup" and dest.exists():
                dest.write_text(dest.read_text(encoding="utf-8").rstrip() + "\n\n" + md, encoding="utf-8")
            else:
                dest.write_text(md, encoding="utf-8")

            (res.updated if prior else res.added).append(relnote)
            files[rel] = {"hash": h, "mtime": st.st_mtime,
                          "size": st.st_size, "note": relnote}
            written += 1

            if want_deadlines and len(res.deadlines) < 20:
                res.deadlines.extend(detect_deadlines(text, parsed))
                if len(res.deadlines) > 20:
                    del res.deadlines[20:]
        except Exception as e:  # per-file resilience (spec §9): record, cap, continue
            if len(res.errors) < 20:
                res.errors.append(f"{rel}: {e}")
            continue

    # Skip the prune when the walk yielded nothing but we have history: an
    # unmounted drive / placeholder dir / permission error must not be read as
    # a mass deletion (spec §4.3).
    if seen_rel or not files:
        for gone in [r for r in files if r not in seen_rel]:
            files.pop(gone, None)

    write_moc(vault, manifest, root, mode)
    return res


def _write_course_hub(vault: Path, course_l: str, notes: list[str]) -> None:
    """(Re)write coursework/<course>.md -- a real note for the [[<course>]]
    link every per-file note carries to resolve against. Without this, every
    per-file note links to a note that never exists and the vault graph shows
    nothing but orphans. Regenerated in full every run from the current
    manifest, same as _moc.md."""
    lines = [
        "---",
        "type: coursework-hub",
        f"course: {course_l.upper()}",
        "---",
        "",
        f"# {course_l.upper()}",
        "",
    ]
    for note in notes:
        target = note[:-3] if note.endswith(".md") else note
        lines.append(f"- [[{target}]]")
    (vault / "coursework" / f"{course_l}.md").write_text(
        "\n".join(lines) + "\n", encoding="utf-8"
    )


def write_moc(vault: Path, manifest: dict, root: Path | str = "D:\\University",
              mode: str = "per-file") -> None:
    files = manifest.get("files", {})
    by_course: dict[str, list[str]] = {}
    for meta in files.values():
        course = meta["note"].split("/")[1].split(".")[0]
        by_course.setdefault(course, []).append(meta["note"])

    lines = ["# Coursework — Map of Content", ""]
    for course in sorted(by_course):
        lines.append(f"- [[{course}]] — {len(by_course[course])} note(s)")

    known_notes = {m["note"] for m in files.values()}
    if mode == "per-file":
        # Hub notes aren't manifest entries -- they're regenerated fresh
        # every run -- so without this every one of them reads as orphaned.
        known_notes |= {f"coursework/{course}.md" for course in by_course}
    cw = vault / "coursework"
    disk_notes = sorted(
        q.relative_to(vault).as_posix() for q in cw.rglob("*.md")
        if q.name != "_moc.md"
    ) if cw.is_dir() else []
    orphans = [n for n in disk_notes if n not in known_notes]
    if orphans:
        lines += ["", "## Orphaned", "",
                  f"_Source file removed from {root}; note kept._", ""]
        lines += [f"- {n}" for n in orphans]

    cw.mkdir(parents=True, exist_ok=True)
    (cw / "_moc.md").write_text("\n".join(lines) + "\n", encoding="utf-8")

    if mode == "per-file":
        for course, notes in by_course.items():
            _write_course_hub(vault, course, sorted(notes))
